import os
import subprocess
import platform
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from agents.expansion_agent import expand_outline
from preset_slides import get_short_presentation_slides  # New module for fallback content

# 🧠 Expand slide content using your LLM agent
def expand_slide(section_title, summary, mode="short"):
    return expand_outline([section_title], summary, mode=mode)[0]

# 🎨 Create a styled slide with custom title and content formatting
def create_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # 🏷️ Title box
    title_box = slide.shapes.add_textbox(Pt(50), Pt(30), Pt(860), Pt(80))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 70, 140)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # 📝 Content box
    content_box = slide.shapes.add_textbox(Pt(50), Pt(130), Pt(860), Pt(400))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    lines = content if isinstance(content, list) else content.split("\n")
    for line in lines:
        if line.strip():
            p = content_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(50, 50, 50)
            p.alignment = PP_ALIGN.LEFT

    return slide

# 🧠 Build the full presentation from outline and summary
def build_presentation(outline, summary, mode="short"):
    prs = Presentation()

    # Use preset slides if outline is ["PRESET"]
    if mode == "short" and outline == ["PRESET"]:
        for slide in get_short_presentation_slides():
            create_slide(prs, slide["title"], slide["content"])
    else:
        for section in outline:
            content = expand_slide(section, summary, mode=mode)
            create_slide(prs, section, content)

    return prs

# 💾 Save and auto-open the presentation cross-platform
def save_presentation(prs, filename="output/generated_presentation.pptx"):
    os.makedirs(os.path.dirname(filename), exist_ok=True)
    prs.save(filename)
    print(f"✅ Presentation saved as {filename}")

    filepath = os.path.abspath(filename)
    open_file(filepath)

# 🚀 Cross-platform file opener
def open_file(filepath):
    system = platform.system()
    try:
        if system == "Darwin":
            subprocess.run(["open", filepath])
        elif system == "Windows":
            os.startfile(filepath)
        elif system == "Linux":
            subprocess.run(["xdg-open", filepath])
    except Exception as e:
        print(f"⚠️ Could not open file automatically: {e}")
